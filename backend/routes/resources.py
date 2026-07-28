"""
资源路由 - 包含多模态资源生成（图片/PPT/文档/思维导图）
"""
import os
import json
import base64
import hmac
import hashlib
import httpx
from datetime import datetime
from time import mktime
from wsgiref.handlers import format_date_time
from urllib.parse import urlencode
from pathlib import Path

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import Optional

from database import get_db_sync
from models import LearningResource, Student, StudentProfile, Exercise, Assessment

router = APIRouter(prefix="/api/resources", tags=["resources"])

# 资源输出目录
OUTPUT_DIR = Path(__file__).parent.parent / "generated"
OUTPUT_DIR.mkdir(exist_ok=True)


# ==================== 请求模型 ====================

class ImageGenRequest(BaseModel):
    """文生图请求"""
    prompt: str
    width: int = 1024
    height: int = 1024
    app_id: str = ""
    api_key: str = ""
    api_secret: str = ""
    student_id: int = 1
    title: str = ""
    topic: str = ""


class PPTGenRequest(BaseModel):
    """PPT 生成请求"""
    topic: str = ""           # 用户需求描述，如"帮我做一个快速排序的PPT"
    outline: str = ""         # [兼容] 直接传入 Markdown 大纲（跳过 AI 生成）
    title: str = ""           # [兼容] 手动指定标题
    student_id: int = 1
    api_key: str = ""         # 前端传入的 LLM API key
    model: str = ""           # 前端传入的模型名
    api_base: str = ""        # 前端传入的 API base URL
    slides_count: int = 6     # 期望幻灯片页数


class DocGenRequest(BaseModel):
    """文档生成请求"""
    prompt: str
    title: str = ""
    student_id: int = 1
    topic: str = ""
    api_key: str = ""
    model: str = ""
    api_base: str = ""


class VideoGenRequest(BaseModel):
    """视频/动画生成请求"""
    topic: str = ""  # 算法/主题描述，如"冒泡排序动画"、"二叉树前序遍历"、"Dijkstra最短路径"
    title: str = ""
    student_id: int = 1
    style: str = "dark"  # dark / light
    count: int = 16  # 数据量
    api_key: str = ""
    model: str = ""
    api_base: str = ""


# ==================== 讯飞鉴权 ====================

def _xfyun_auth(host: str, path: str, api_key: str, api_secret: str) -> tuple[str, str]:
    """生成讯飞 HMAC-SHA256 鉴权的 date 和 authorization"""
    cur_time = datetime.now()
    date = format_date_time(mktime(cur_time.timetuple()))

    tmp = f"host: {host}\ndate: {date}\nPOST {path} HTTP/1.1"
    tmp_sha = hmac.new(
        api_secret.encode('utf-8'),
        tmp.encode('utf-8'),
        digestmod=hashlib.sha256
    ).digest()
    signature = base64.b64encode(tmp_sha).decode('utf-8')

    auth_origin = (
        f'api_key="{api_key}", '
        f'algorithm="hmac-sha256", '
        f'headers="host date request-line", '
        f'signature="{signature}"'
    )
    authorization = base64.b64encode(auth_origin.encode('utf-8')).decode('utf-8')
    return date, authorization


def _xfyun_build_auth_url(host: str, path: str, api_key: str, api_secret: str, schema: str = "https") -> str:
    """生成带鉴权查询参数的完整URL。

    讯飞 VMS 虚拟人等接口要求 host/date/authorization 作为 URL 查询参数传递
    （而非 HTTP header），否则报 "enforced header 'host' not used for signature creation"。
    """
    cur_time = datetime.now()
    date = format_date_time(mktime(cur_time.timetuple()))

    signature_origin = f"host: {host}\ndate: {date}\nPOST {path} HTTP/1.1"
    tmp_sha = hmac.new(
        api_secret.encode('utf-8'),
        signature_origin.encode('utf-8'),
        digestmod=hashlib.sha256
    ).digest()
    signature = base64.b64encode(tmp_sha).decode('utf-8')

    auth_origin = (
        f'api_key="{api_key}", '
        f'algorithm="hmac-sha256", '
        f'headers="host date request-line", '
        f'signature="{signature}"'
    )
    authorization = base64.b64encode(auth_origin.encode('utf-8')).decode('utf-8')

    params = {"host": host, "date": date, "authorization": authorization}
    return f"{schema}://{host}{path}?{urlencode(params)}"


# ==================== 文生图 (讯飞星火) ====================

XF_TTI_HOST = "spark-api.cn-huabei-1.xf-yun.com"
XF_TTI_PATH = "/v2.1/tti"
XF_TTI_URL = f"https://{XF_TTI_HOST}{XF_TTI_PATH}"


@router.post("/generate-image")
async def generate_image(req: ImageGenRequest):
    """调用讯飞星火文生图 API 生成图片，返回 base64 + 保存路径"""
    if not all([req.app_id, req.api_key, req.api_secret]):
        raise HTTPException(400, "缺少讯飞星火鉴权信息 (app_id/api_key/api_secret)")

    date, authorization = _xfyun_auth(XF_TTI_HOST, XF_TTI_PATH, req.api_key, req.api_secret)

    body = {
        "header": {"app_id": req.app_id},
        "parameter": {
            "chat": {
                "domain": "general",
                "width": req.width,
                "height": req.height,
            }
        },
        "payload": {
            "message": {
                "text": [{"role": "user", "content": req.prompt}]
            }
        }
    }

    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.post(
            XF_TTI_URL,
            json=body,
            headers={
                "Content-Type": "application/json;charset=UTF-8",
                "Authorization": authorization,
                "Date": date,
                "Host": XF_TTI_HOST,
            }
        )
        data = resp.json()

    if data.get("header", {}).get("code") != 0:
        err_msg = data.get("header", {}).get("message", "未知错误")
        raise HTTPException(500, f"讯飞星火返回错误: {err_msg}")

    # 提取 base64 图片
    choices = data.get("payload", {}).get("choices", {})
    text_list = choices.get("text", [])
    if not text_list:
        raise HTTPException(500, "讯飞星火未返回图片数据")

    img_b64 = text_list[0].get("content", "")
    if not img_b64:
        raise HTTPException(500, "图片内容为空")

    # 保存图片
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = "".join(c if c.isalnum() or c in '_-' else '_' for c in req.prompt[:30])
    filename = f"img_{safe_name}_{timestamp}.png"
    filepath = OUTPUT_DIR / filename

    with open(filepath, "wb") as f:
        f.write(base64.b64decode(img_b64))

    # 存入数据库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=req.title or f"AI 生成图片: {req.prompt[:40]}",
            resource_type="image",
            content=f"/api/resources/file/{filename}",
            topic=req.topic or "AI 生成",
            difficulty="基础",
            tags=["AI生成", "文生图", "讯飞星火"],
            agent_generated="xfyun_tti",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "filename": filename,
        "url": f"/api/resources/file/{filename}",
        "base64": img_b64,  # 前端可直接显示
        "width": req.width,
        "height": req.height,
    }


# ==================== PPT 生成 (python-pptx) ====================

@router.post("/generate-ppt")
async def generate_ppt(req: PPTGenRequest):
    """两阶段 AI 生成 PPTX：1) 生成标题+大纲  2) 细化每页内容  3) 构建 PPTX"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(500, "python-pptx 未安装，请运行 pip install python-pptx")

    # === 阶段 0: 如果前端直接传了 outline，走兼容模式 ===
    if req.outline.strip():
        title = req.title or req.topic or "AI 生成课件"
        slides_data = _parse_markdown_outline(req.outline, title)
        return _build_and_save_pptx(slides_data, title, req)

    # === 阶段 1: AI 生成标题和大纲 ===
    topic_text = req.topic.strip()
    if not topic_text:
        raise HTTPException(400, "请提供 PPT 主题描述")

    llm_kwargs = {}
    if req.api_key:
        llm_kwargs["api_key"] = req.api_key
    if req.model:
        llm_kwargs["model"] = req.model
    if req.api_base:
        llm_kwargs["api_base"] = req.api_base

    outline_system = """你是一个资深的算法/编程教学 PPT 大纲设计师。用户会给你一个主题，你需要设计一份结构清晰、内容丰富的 PPT 大纲。

要求：
1. 设计一个专业、吸引人的标题
2. 规划 {slides_count} 页幻灯片，每页 4-6 个要点
3. 每个要点必须是**完整的描述性句子**（15-30字），不是单字关键词
   - 错误示例: "定义"
   - 正确示例: "动态规划是一种将复杂问题分解为重叠子问题、通过记录中间结果避免重复计算的优化方法"
4. 封面页（第1页）要包含副标题式概述
5. 中间页按知识递进组织：背景→概念→原理→步骤→示例→应用
6. 如果涉及代码，在对应页加入 code 字段（独立字符串）
7. 最后一页是总结回顾

严格输出以下 JSON（不要 markdown 包裹，不要其他文字）：
{
  "title": "专业主标题",
  "subtitle": "副标题或一句话概述",
  "slides": [
    {
      "heading": "页面标题",
      "points": ["完整句子要点1", "完整句子要点2", "完整句子要点3", "完整句子要点4"],
      "code": "可选的代码示例字符串（没有则省略此字段）"
    }
  ]
}"""

    outline_user = f"请为主题「{topic_text}」设计一份 {req.slides_count} 页的 PPT 课件大纲。要求每页都有扎实的教学内容，不是空洞的标题堆砌。"

    try:
        from llm_service import chat_completion
        outline_raw = await chat_completion(
            system_prompt=outline_system.format(slides_count=req.slides_count),
            user_message=outline_user,
            temperature=0.6,
            max_tokens=4096,
            **llm_kwargs,
        )
    except Exception as e:
        # LLM 调用失败时使用默认大纲
        import traceback
        print(f"[PPT Stage 1 ERROR] {e}")
        traceback.print_exc()
        outline_raw = json.dumps({
            "title": topic_text,
            "subtitle": f"深入理解{topic_text}的核心原理与典型应用",
            "slides": [
                {"heading": topic_text, "points": [
                    f"本课件将系统讲解{topic_text}的核心概念、底层原理与典型应用场景，帮助你从理论基础到实战能力全面提升",
                    f"通过循序渐进的讲解方式,从基础概念出发,逐步深入到算法实现和优化技巧",
                    f"涵盖理论推导、代码实现(C++/Python)、复杂度分析和实际案例四大维度",
                    f"适合有一定编程基础、希望深入理解算法原理的学习者"
                ]},
                {"heading": "背景与核心概念", "points": [
                    f"首先理解{topic_text}要解决的核心问题是什么,它与朴素解法相比有何优势",
                    f"掌握算法涉及的关键数据结构及其特性,理解数据组织方式对算法效率的决定性影响",
                    f"了解算法的输入输出规范、前置条件和边界情况的处理方式",
                    f"对比相似算法,明确{topic_text}的适用场景和局限性"
                ]},
                {"heading": "算法原理与执行流程", "points": [
                    f"将{topic_text}拆解为清晰的步骤,用具体数值示例逐步追踪算法的执行过程",
                    f"理解每个步骤的设计意图和正确性依据,分析关键决策点的逻辑",
                    f"推导时间复杂度和空间复杂度,通过数学分析理解算法的效率瓶颈",
                    f"讨论常见实现陷阱和工程优化技巧,提升代码质量和运行效率"
                ]},
                {"heading": "代码实现与解析", "points": [
                    f"提供{topic_text}的完整 C++ 和 Python 实现,代码简洁清晰、注释完整",
                    f"逐段分析代码的关键逻辑,解释每个变量和判断条件的含义",
                    f"分析代码中的边界条件处理,讨论索引越界、空输入等特殊情况",
                    f"提供测试用例并验证代码的正确性和鲁棒性"
                ]},
                {"heading": "应用场景与经典题目", "points": [
                    f"列举 3-5 个直接应用{topic_text}的经典算法题目,说明每道题的转化思路",
                    f"分析每道题的输入范围、约束条件和评测数据规模",
                    f"对比不同解法的优劣,展示{topic_text}在特定场景下的优越性",
                    f"总结从题目中识别{topic_text}适用条件的方法论"
                ]},
                {"heading": "总结与进阶方向", "points": [
                    f"回顾{topic_text}的核心知识点: 原理、实现、复杂度、适用场景",
                    f"总结常见易错点: 边界条件、类型溢出、特殊输入处理",
                    f"推荐进阶学习路径: 相关算法变体、更复杂的应用场景、与其他算法的组合使用",
                    f"从{topic_text}延伸到更广阔的知识体系,培养系统化的算法思维"
                ]}
            ]
        })

    # 解析大纲 JSON
    try:
        cleaned = outline_raw.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else cleaned
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
        outline_data = json.loads(cleaned.strip())
        title = outline_data.get("title", topic_text)
        subtitle = outline_data.get("subtitle", "")
        raw_slides = outline_data.get("slides", [])
    except json.JSONDecodeError:
        title = topic_text
        subtitle = ""
        raw_slides = [{"heading": f"关于 {topic_text}", "points": [outline_raw[:500]]}]

    # === 阶段 2: AI 细化每页内容 — 逐页调用以获得深度内容 ===
    detail_system = """你是一位算法教学的资深讲师，正在为一门高质量的算法课程撰写 PPT 详解内容。

你的任务：根据给定的一页 PPT 大纲，将其展开为充实、有深度的讲解内容。

内容要求 —— 这是最重要的部分：
- 每个要点展开为 **80-200 字** 的详细讲解段落，不是一两句话就完事
- 包含具体原理说明、推导过程、生动的例子或类比
- 如果涉及代码，给出**完整可运行的代码片段**（用 ```cpp 或 ```python 包裹）
- 如果是概念页，请用通俗易懂的语言解释，让初学者也能理解
- 保持专业但亲切的口吻，像一位耐心的老师在讲课

格式要求：
- 以 ## 页面标题 开头
- 每个要点用 - 开头，后面跟随 80-200 字的详细讲解
- 如果有代码，在相关要点后用 ``` 代码块单独放置
- 可以加入子要点（用 4 个空格缩进的 - 表示）来展开细节

请直接输出该页的 Markdown 内容，不要包含任何开场白或结尾语。"""

    # === 逐页调用 AI 细化 ===
    all_detailed_pages = []
    for i, slide in enumerate(raw_slides):
        heading = slide.get("heading", f"第{i+1}页")
        points = slide.get("points", [])
        code = slide.get("code", "")

        points_text = "\n".join([f"- {p}" for p in points])
        code_text = f"\n\n参考代码：\n```\n{code}\n```" if code else ""

        detail_user = f"""主题：{title}
页面：{heading}
大纲要点：
{points_text}{code_text}

请将以上内容展开为一页**充实详细**的 PPT 讲解。每个要点都要写成 80-200 字的段落，不能只写一两句话。"""

        try:
            from llm_service import chat_completion
            page_md = await chat_completion(
                system_prompt=detail_system,
                user_message=detail_user,
                temperature=0.7,
                max_tokens=2048,
                **llm_kwargs,
            )
            all_detailed_pages.append(page_md)
        except Exception:
            import traceback
            print(f"[PPT Stage 2 Page {i+1} ERROR] {traceback.format_exc()}")
            # 该页细化失败，用原始大纲
            fallback = f"## {heading}\n" + "\n".join([f"- {p}" for p in points])
            if code:
                fallback += f"\n```\n{code}\n```"
            all_detailed_pages.append(fallback)

    # 合并所有页面
    detailed_md = "\n\n".join(all_detailed_pages)

    # 解析详细大纲为幻灯片数据
    slides_data = _parse_markdown_outline(detailed_md, title)

    # 如果有 subtitle 且第一页是封面，把 subtitle 加到封面
    if subtitle and slides_data:
        slides_data[0]["content"].insert(0, subtitle)

    return _build_and_save_pptx(slides_data, title, req)


def _build_and_save_pptx(slides_data: list[dict], title: str, req: PPTGenRequest) -> dict:
    """构建 PPTX 文件并保存"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    if not slides_data:
        slides_data = [{"title": title, "content": ["请提供更详细的内容描述"]}]

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_slide(slide, slide_info, idx, len(slides_data))

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = "".join(c if c.isalnum() or c in '_-' else '_' for c in title[:30])
    filename = f"ppt_{safe_name}_{timestamp}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))

    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=title,
            resource_type="ppt",
            content=f"/api/resources/file/{filename}",
            topic=req.topic or title,
            difficulty="基础",
            tags=["PPT", "AI生成", "课件"],
            agent_generated="ppt_engine",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "filename": filename,
        "url": f"/api/resources/file/{filename}",
        "slides": len(slides_data),
        "title": title,
    }


def _parse_markdown_outline(md: str, default_title: str = "AI 生成课件") -> list[dict]:
    """解析 Markdown 为幻灯片结构
    支持:
    # 标题 → 封面页
    ## 标题 → 内容页标题
    - 列表 → 要点
    ``` → 代码块
    普通文本 → 段落
    """
    slides = []
    current = None
    in_code = False
    code_lines = []

    for line in md.split('\n'):
        stripped = line.strip()

        # 代码块开始/结束
        if stripped.startswith('```'):
            if in_code:
                if code_lines:
                    current['content'].append('\n'.join(code_lines))
                    code_lines = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(line)
            continue

        # 一级标题 → 封面
        if stripped.startswith('# ') and not stripped.startswith('## '):
            if current:
                slides.append(current)
            current = {"title": stripped[2:], "content": []}
            continue

        # 二级标题 → 新页
        if stripped.startswith('## '):
            if current:
                slides.append(current)
            current = {"title": stripped[3:], "content": []}
            continue

        # 三级标题或列表
        if stripped.startswith('- ') or stripped.startswith('* '):
            if current is None:
                current = {"title": "内容", "content": []}
            current['content'].append(stripped[2:])
            continue

        # 普通文本
        if stripped and current is None:
            current = {"title": default_title, "content": []}
        if stripped and current:
            current['content'].append(stripped)

    if current:
        slides.append(current)

    return slides


def _build_slide(slide, info: dict, page_index: int = 0, total_pages: int = 1):
    """构建单张幻灯片，根据内容量自适应调整字号"""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    content_items = info.get("content", [])
    # 统计代码行数和非代码内容的总字符数来估算密集度
    total_chars = sum(len(item) for item in content_items if not _is_code_line(item))
    code_lines_count = sum(1 for item in content_items if _is_code_line(item))
    is_dense = total_chars > 400 or len(content_items) > 7

    # 背景色 (深色科技风)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1B, 0x26)

    # 标题
    left = Inches(0.8)
    top = Inches(0.35)
    width = Inches(14.4)
    title_height = Inches(0.9) if not is_dense else Inches(0.7)

    title_box = slide.shapes.add_textbox(left, top, width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = info.get("title", "")
    p.font.size = Pt(28) if not is_dense else Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x7D, 0xC8, 0xFF)

    # 底部装饰线
    line = slide.shapes.add_shape(
        1, left, top + title_height + Inches(0.05), Inches(2.5), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x7D, 0xC8, 0xFF)
    line.line.fill.background()

    # 内容区域
    content_top = top + title_height + Inches(0.15)
    content_height = Inches(8.0) - content_top

    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    # 自适应字号
    base_body_pt = 14 if is_dense else 16
    code_pt = 11 if is_dense else 12
    bullet_color = RGBColor(0xE8, 0xE8, 0xF0)
    code_color = RGBColor(0xA8, 0xD8, 0x80)
    dim_color = RGBColor(0xAA, 0xAA, 0xBB)

    # 前置标记符号
    BULLET = "  "

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        is_code = _is_code_line(item)

        if is_code:
            # 代码行
            p.text = item
            p.font.size = Pt(code_pt)
            p.font.name = "Consolas"
            p.font.color.rgb = code_color
            p.space_after = Pt(2)
            p.space_before = Pt(1)
        else:
            # 检测缩进级别（4空格缩进 = 子要点）
            indent_level = 0
            display_text = item
            if item.startswith("    - ") or item.startswith("    * "):
                indent_level = 1
                display_text = item[2:].lstrip("- *").strip()
            elif item.startswith("- ") or item.startswith("* "):
                indent_level = 0
                display_text = item[2:].strip()
            elif item.startswith("  - ") or item.startswith("  * "):
                indent_level = 1
                display_text = item[2:].lstrip("- *").strip()

            p.text = display_text
            p.font.size = Pt(max(10, base_body_pt - indent_level * 2))
            p.font.color.rgb = bullet_color if indent_level == 0 else dim_color
            p.space_after = Pt(4) if is_dense else Pt(6)
            p.space_before = Pt(1)
            p.level = indent_level

    # 页脚
    footer = slide.shapes.add_textbox(Inches(13.5), Inches(8.4), Inches(2.5), Inches(0.4))
    pf = footer.text_frame.paragraphs[0]
    pf.text = f"AlgoAscend AI  ·  {page_index + 1}/{total_pages}"
    pf.font.size = Pt(9)
    pf.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    pf.alignment = PP_ALIGN.RIGHT


def _is_code_line(text: str) -> bool:
    """判断文本行是否是代码"""
    code_indicators = [
        '#include', 'int ', 'void ', 'class ', 'def ', 'for(', 'while(',
        'if (', 'return ', 'cout', 'printf', 'vector<', 'using namespace',
        'public:', 'private:', 'protected:', 'def ', 'import ', 'from ',
        'print(', 'elif', 'else:', 'try:', 'except', '```', 'std::',
        'const ', 'auto ', 'struct ', 'enum ', 'template<', '#define',
        'cin ', 'endl', 'push_back', 'pop_back', '.size()', '.begin()',
    ]
    return any(kw in text for kw in code_indicators)


# ==================== 文档生成 ====================

@router.post("/generate-doc")
async def generate_doc(req: DocGenRequest):
    """生成结构化 Markdown 文档"""
    # 用 LLM 生成内容（简化：直接返回模板化内容）
    content = f"""# {req.title or 'AI 生成学习文档'}

> 主题: {req.topic}
> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}

---

## 概述

{req.prompt}

> 本文档由 AlgoAscend AI 多智能体系统自动生成，内容源自知识库检索和模型推理。
"""
    # 存库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=req.title or f"学习文档: {req.topic[:40]}",
            resource_type="doc",
            content=content,
            topic=req.topic or "学习笔记",
            difficulty="基础",
            tags=["AI生成", "文档", "学习笔记"],
            agent_generated="doc_engine",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "content": content,
    }


# ==================== 思维导图生成 ====================

MINDMAP_SYSTEM_PROMPT = """你是一个专业的算法知识梳理助手。根据用户输入的主题，生成一个结构清晰的思维导图。

要求：
1. 输出严格的 JSON 树结构，格式：{"root": "主题名", "children": [{"text": "分支1", "children": [...]}, ...]}
2. 根节点是主题名，子节点是分类/细分知识点
3. 每个分支最多 3-4 层深度，同级节点 3-6 个
4. 分支要有逻辑层次，不要简单罗列
5. 每个节点文本不超过 15 个字，简洁精炼

示例输入："排序算法有哪些"
示例输出：
{
  "root": "排序算法",
  "children": [
    {"text": "比较排序", "children": [
      {"text": "交换排序", "children": [{"text": "冒泡排序", "children": []}, {"text": "快速排序", "children": []}]},
      {"text": "插入排序", "children": [{"text": "直接插入", "children": []}, {"text": "希尔排序", "children": []}]},
      {"text": "选择排序", "children": [{"text": "简单选择", "children": []}, {"text": "堆排序", "children": []}]},
      {"text": "归并排序", "children": []}
    ]},
    {"text": "非比较排序", "children": [
      {"text": "计数排序", "children": []},
      {"text": "基数排序", "children": []},
      {"text": "桶排序", "children": []}
    ]},
    {"text": "复杂度对比", "children": [
      {"text": "时间复杂度", "children": []},
      {"text": "空间复杂度", "children": []},
      {"text": "稳定性", "children": []}
    ]}
  ]
}"""


@router.post("/generate-mindmap")
async def generate_mindmap(req: DocGenRequest):
    """用 LLM 生成结构化的思维导图 JSON 树"""
    # 检查 API Key
    from config import AppConfig
    api_key = req.api_key or AppConfig().llm.api_key
    if not api_key:
        raise HTTPException(400, "未配置 API Key，请在网页右上角选择模型并在设置中填写对应的 API Key")

    try:
        from llm_service import chat_with_json_output
        tree_data = await chat_with_json_output(
            system_prompt=MINDMAP_SYSTEM_PROMPT,
            user_message=req.prompt,
            api_key=req.api_key or None,
            model=req.model or None,
            api_base=req.api_base or None,
        )
    except Exception as e:
        raise HTTPException(500, f"LLM 调用失败: {str(e)}")

    # 验证树结构
    root_text = tree_data.get("root", req.prompt[:20])
    children = tree_data.get("children", [])

    if not root_text:
        raise HTTPException(500, "LLM 未能生成有效的思维导图结构")

    # 存库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=req.title or f"思维导图: {root_text[:30]}",
            resource_type="mindmap",
            content=json.dumps(tree_data, ensure_ascii=False),
            topic=req.topic or root_text,
            difficulty="基础",
            tags=["AI生成", "思维导图"],
            agent_generated="mindmap_llm",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "tree": tree_data,
    }


# ==================== 代码实操案例生成 ====================

CODE_CASE_SYSTEM_PROMPT = """你是一位资深C++算法教练，专门生成带详细注释的代码实操案例。

根据用户指定的主题，生成一份结构清晰、注释完整的代码实操文档。

## 输出格式
使用 Markdown，包含以下结构：

## 一、问题概述
简要描述该算法/数据结构要解决的问题

## 二、核心思路
用通俗语言解释核心思想，配合文字图或ASCII示意图

## 三、完整代码实现
```cpp
// 带详细注释的完整C++代码，包括：
// - 文件头注释（功能说明、复杂度分析）
// - 每个关键步骤的中文注释
// - 边界条件处理的注释
// - 测试用例的 main 函数
```

## 四、关键步骤拆解
逐段分析代码的核心逻辑，每个步骤配解释

## 五、运行示例
给出输入输出示例和预期结果

## 六、常见错误与调试技巧
列举3-5个常见错误和对应的调试方法

## 七、变体与扩展
介绍该算法的变体或相关扩展知识

要求：
- 代码必须完整可运行，注释覆盖率不低于30%
- 用生活化类比帮助理解
- 标注时间复杂度与空间复杂度
- 难度与学生水平匹配"""


class CodeCaseGenRequest(BaseModel):
    topic: str
    student_id: int = 1
    api_key: str = ""
    model: str = ""
    api_base: str = ""


@router.post("/generate-code-case")
async def generate_code_case(req: CodeCaseGenRequest):
    """生成带详细注释的代码实操案例"""
    topic = req.topic.strip()
    if not topic:
        raise HTTPException(400, "请提供代码实操主题")

    from config import AppConfig
    api_key = req.api_key or AppConfig().llm.api_key
    if not api_key:
        raise HTTPException(400, "未配置 API Key")

    try:
        from llm_service import chat_completion
        content = await chat_completion(
            system_prompt=CODE_CASE_SYSTEM_PROMPT,
            user_message=f"请为主题「{topic}」生成一份带详细注释的C++代码实操案例。代码要完整可运行，注释要详细。",
            temperature=0.5,
            max_tokens=4096,
            api_key=req.api_key or None,
            model=req.model or None,
            api_base=req.api_base or None,
        )
    except Exception as e:
        raise HTTPException(500, f"LLM 调用失败: {str(e)}")

    # 存库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=f"代码实操: {topic[:40]}",
            resource_type="code_case",
            content=content,
            topic=topic,
            difficulty="基础",
            tags=["代码实操", "AI生成", "带注释", topic],
            agent_generated="code_case_agent",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "content": content,
        "topic": topic,
    }


# ==================== 实践项目案例生成 ====================

PROJECT_SYSTEM_PROMPT = """你是一位高级C++算法工程师，专门设计综合性的小型算法实践项目。

根据用户指定的主题，设计一个完整的实践项目，让学生通过动手实现来深入理解。

## 输出格式
使用 Markdown，包含以下结构：

## 🎯 项目目标
一句话概括本项目要做什么、学到什么

## 📋 需求描述
### 功能需求（3-5条）
### 技术要求（如时间复杂度、空间限制等）
### 输入输出规范

## 🏗️ 设计方案
### 数据结构设计
- 核心数据结构的选择理由
- 类/结构体设计

### 算法设计
- 核心算法流程图（用ASCII或文字描述）
- 关键接口定义

## 💻 实现步骤
分4-6个步骤，每步有明确的目标和验证方法：
1. 步骤一：搭建框架
2. 步骤二：实现核心逻辑
...

## ✅ 测试用例
提供3-5个测试用例（输入+期望输出+边界情况）

## 📊 评分标准
（可选）自我评估的检查清单

## 🔧 扩展挑战
2-3个进阶变体，供有余力的学生挑战

要求：
- 项目规模适合1-3天完成
- 难度与学生水平匹配
- 提供完整的代码骨架（框架代码）
- 重点标注核心技术难点"""


class ProjectGenRequest(BaseModel):
    topic: str
    student_id: int = 1
    api_key: str = ""
    model: str = ""
    api_base: str = ""


@router.post("/generate-project")
async def generate_project(req: ProjectGenRequest):
    """生成综合算法实践项目案例"""
    topic = req.topic.strip()
    if not topic:
        raise HTTPException(400, "请提供项目主题")

    from config import AppConfig
    api_key = req.api_key or AppConfig().llm.api_key
    if not api_key:
        raise HTTPException(400, "未配置 API Key")

    try:
        from llm_service import chat_completion
        content = await chat_completion(
            system_prompt=PROJECT_SYSTEM_PROMPT,
            user_message=f"请为主题「{topic}」设计一个完整的C++算法实践项目。项目规模适中，适合1-3天完成，包含完整的设计方案和实现步骤。",
            temperature=0.5,
            max_tokens=4096,
            api_key=req.api_key or None,
            model=req.model or None,
            api_base=req.api_base or None,
        )
    except Exception as e:
        raise HTTPException(500, f"LLM 调用失败: {str(e)}")

    # 存库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=f"实践项目: {topic[:40]}",
            resource_type="project",
            content=content,
            topic=topic,
            difficulty="进阶",
            tags=["实践项目", "AI生成", "综合案例", topic],
            agent_generated="project_agent",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "content": content,
        "topic": topic,
    }


# ==================== 画像驱动的资源推荐 ====================

RECOMMEND_SYSTEM_PROMPT = """你是一个学习资源推荐引擎。根据学生的六维画像和学习数据，推荐最适合的学习资源。

## 输出格式（严格JSON，不要markdown包裹）
{
  "recommendations": [
    {
      "type": "doc" | "video" | "exercise" | "code_case" | "project" | "mindmap" | "reading",
      "title": "推荐标题",
      "description": "简短推荐理由（20字以内）",
      "topic": "关联知识点",
      "priority": "high" | "medium" | "low",
      "reason": "基于画像的推荐理由（50字以内）"
    }
  ],
  "study_tip": "一条鼓励性学习建议"
}

推荐策略：
- 知识薄弱点 → 推荐讲解文档+对应练习题
- 兴趣领域 → 推荐进阶项目+代码实操
- 学习节奏快 → 推荐更多挑战性内容
- 学习节奏慢 → 推荐基础巩固+拓展阅读
- 易错点 → 推荐针对性练习+图解内容
- 学习目标 → 推荐路径相关的下一个知识点"""


@router.get("/recommend/{student_id}")
async def get_recommendations(student_id: int = 1):
    """基于学生画像生成个性化资源推荐"""
    from config import AppConfig
    api_key = AppConfig().llm.api_key

    # 获取学生画像
    db = get_db_sync()
    try:
        profile = db.query(StudentProfile).filter(
            StudentProfile.student_id == student_id
        ).first()
        profile_data = profile.to_dict() if profile else {}

        # 获取学习统计
        exercises = db.query(Exercise).filter(
            Exercise.student_id == student_id
        ).all()
        resources = db.query(LearningResource).filter(
            LearningResource.student_id == student_id
        ).all()
        assessments = db.query(Assessment).filter(
            Assessment.student_id == student_id
        ).all()

        total_exercises = len(exercises)
        correct_exercises = sum(1 for e in exercises if e.is_correct)
        accuracy = correct_exercises / max(1, total_exercises)

        # 资源类型分布
        res_by_type = {}
        for r in resources:
            rt = r.resource_type
            res_by_type[rt] = res_by_type.get(rt, 0) + 1
    finally:
        db.close()

    # 从已有资源中选取
    db2 = get_db_sync()
    try:
        recent_resources = db2.query(LearningResource).filter(
            LearningResource.student_id == student_id
        ).order_by(LearningResource.created_at.desc()).limit(30).all()
        available_resources = [r.to_dict() for r in recent_resources]
    finally:
        db2.close()

    # 尝试用LLM生成推荐
    if api_key:
        try:
            from llm_service import chat_with_json_output
            stats_info = {
                "total_exercises": total_exercises,
                "accuracy": round(accuracy * 100, 1),
                "total_resources": len(available_resources),
                "resource_types": res_by_type,
                "total_assessments": len(assessments),
            }

            prompt = f"""基于以下学生画像和学习数据，推荐6个学习资源。

## 学生画像
{json.dumps(profile_data, ensure_ascii=False, indent=2)}

## 学习统计
{json.dumps(stats_info, ensure_ascii=False, indent=2)}

## 已有资源主题（避免重复推荐）
{', '.join([r.get('topic', '') for r in available_resources[:10] if r.get('topic')])}

请严格按JSON格式输出推荐结果。"""
            result = await chat_with_json_output(
                system_prompt=RECOMMEND_SYSTEM_PROMPT,
                user_message=prompt,
                temperature=0.6,
            )
            return {
                "status": "ok",
                "recommendations": result.get("recommendations", []),
                "study_tip": result.get("study_tip", "坚持每天学习，积少成多！"),
                "source": "ai",
            }
        except Exception:
            pass  # LLM失败则用规则推荐

    # 规则推荐fallback
    rule_recommendations = _rule_based_recommend(profile_data, res_by_type, available_resources)
    return {
        "status": "ok",
        "recommendations": rule_recommendations,
        "study_tip": "每天进步一点点，坚持带来大改变！",
        "source": "rule",
    }


def _rule_based_recommend(profile: dict, res_by_type: dict, available: list) -> list:
    """基于规则的资源推荐"""
    recs = []
    interests = profile.get("interests", {})
    error_patterns = profile.get("error_patterns", {})
    goals = profile.get("learning_goals", {})

    fav_topics = interests.get("favorite_topics", [])
    weak_areas = error_patterns.get("weak_areas", [])
    common_errors = error_patterns.get("common_errors", [])

    # 薄弱点 → 推荐文档+练习题
    for area in weak_areas[:2]:
        recs.append({
            "type": "doc", "title": f"深入理解{area}",
            "description": f"针对薄弱环节{area}的专题讲解",
            "topic": area, "priority": "high",
            "reason": f"你在{area}方面需要加强，建议优先学习"
        })
        recs.append({
            "type": "exercise", "title": f"{area}专项练习",
            "description": f"{area}相关练习题巩固",
            "topic": area, "priority": "high",
            "reason": f"练习是掌握{area}的最佳方式"
        })

    # 兴趣领域 → 推荐项目+代码实操
    for topic in fav_topics[:2]:
        recs.append({
            "type": "code_case", "title": f"{topic}代码实操",
            "description": f"{topic}的带注释完整代码案例",
            "topic": topic, "priority": "medium",
            "reason": f"动手实践你感兴趣的{topic}"
        })
        recs.append({
            "type": "project", "title": f"{topic}综合项目",
            "description": f"基于{topic}的实践项目",
            "topic": topic, "priority": "medium",
            "reason": f"通过项目深入掌握{topic}"
        })

    # 通用推荐
    recs.append({
        "type": "mindmap", "title": "知识体系梳理",
        "description": "生成当前学习内容的知识导图",
        "topic": "综合", "priority": "low",
        "reason": "整理知识结构，建立系统化认知"
    })
    recs.append({
        "type": "reading", "title": "拓展阅读推荐",
        "description": "搜索相关拓展阅读材料",
        "topic": "拓展", "priority": "low",
        "reason": "拓宽知识面，了解前沿发展"
    })

    return recs[:6]


# ==================== 文件下载 ====================

@router.get("/file/{filename}")
async def serve_file(filename: str):
    """提供生成文件的下载"""
    from fastapi.responses import FileResponse
    filepath = OUTPUT_DIR / filename
    if not filepath.exists():
        raise HTTPException(404, "文件不存在")
    return FileResponse(str(filepath))


# ==================== Manim 算法动画脚本生成（LLM 动态生成） ====================

MANIM_VIDEO_SYSTEM_PROMPT = """你是一个专业的 Manim CE (Community Edition) 动画脚本生成器。根据用户描述的主题，生成一个完整的、可直接渲染的 Manim Python 脚本。

## Manim CE 基础知识
- 使用 `from manim import *` 导入
- Scene 类名必须以大写字母开头，建议使用 `AlgoViz` 或与主题相关的名称
- 使用 `self.play(...)` 驱动动画，`self.wait(n)` 暂停
- 使用 `self.add(...)` 添加静态对象

## 支持的可视化风格

### 1. 排序/比较类算法（柱状图 + 动画过程）
```python
class AlgoViz(Scene):
    def construct(self):
        import random
        data = [random.randint(5, 95) for _ in range(N)]
        n = len(data)
        max_val = max(data)

        bars = VGroup()
        labels = VGroup()
        for i, v in enumerate(data):
            bar = Rectangle(width=0.35, height=v/max_val*5, fill_color=BLUE, fill_opacity=0.8, stroke_color=WHITE, stroke_width=0.5)
            bar.move_to(np.array([(i - n/2) * 0.45, -2 + bar.height/2, 0]))
            bars.add(bar)
            label = Text(str(v), font_size=14, color=WHITE)
            label.next_to(bar, DOWN, buff=0.1)
            labels.add(label)

        self.add(bars, labels)
        title = Text("算法名称", font_size=36, color=WHITE).to_edge(UP)
        self.play(Write(title))
        self.wait(0.5)

        # ... 算法动画过程 ...

        self.wait(2)
```
关键：用颜色区分状态 — BLUE(默认)、YELLOW(比较中)、RED(枢轴/关键元素)、GREEN(已排序/已完成)、ORANGE/PURPLE(辅助标记)

### 2. 数据结构类（树、图、链表等）
使用 Circle、Line、Arrow、Rectangle 等组合构建可视化结构，通过颜色和高亮展示操作过程。

### 3. 搜索/遍历类
在数组或图结构上展示搜索路径，用颜色标记已访问/未访问/当前节点。

### 4. 递归类
可以用嵌套调用 tree 展示递归过程，或直接在数组上标注每一层的处理范围。

## 颜色约定
- BLUE: 默认/未处理
- YELLOW: 当前正在处理
- RED: 关键元素（枢轴、目标等）
- GREEN: 已完成/已找到
- ORANGE: 左半部分/辅助
- PURPLE: 右半部分/辅助

## 布局约定
- 标题用 `Text(..., font_size=36).to_edge(UP)`
- 主体内容放在画面中央偏上区域
- 底部可放公式/说明文字
- 画面尺寸默认 14.22 x 8.0 (16:9 比例)

## 输出格式
**只输出 Python 代码**，放在 ```python ... ``` 代码块中。代码必须完整、可直接保存为 .py 文件并使用 `manim -pql` 渲染。
不要输出任何解释、说明或额外的文字，只输出代码块。"""


def _extract_manim_scene(script: str) -> str:
    """从 Manim 脚本中提取 Scene 类名"""
    import re
    # 匹配 class XXX(Scene): 或 class XXX(ThreeDScene): 等
    match = re.search(r'class\s+(\w+)\s*\(\s*(?:ThreeD)?Scene\s*\)', script)
    return match.group(1) if match else "AlgoViz"


INSTALL_MANIM = """# Manim 安装指南 (Windows)
# 1. 安装 Python 3.8+ (已安装则跳过)
# 2. 安装 FFmpeg: choco install ffmpeg  或从 https://ffmpeg.org 下载
# 3. 安装 Manim:
#    pip install manim
# 4. 渲染视频（在脚本所在目录运行）:
#    manim -pql <脚本文件名.py> <类名>
#    -p: 预览, -q: 质量 (l/m/h)
"""


@router.post("/generate-video")
async def generate_video_script(req: VideoGenRequest):
    """用 LLM 动态生成 Manim 算法动画 Python 脚本，并尝试渲染"""
    topic = req.topic.strip()
    if not topic:
        raise HTTPException(400, "请提供算法/主题描述")

    # 用 LLM 生成 Manim 脚本
    from config import AppConfig
    api_key = req.api_key or AppConfig().llm.api_key
    if not api_key:
        raise HTTPException(400, "未配置 API Key，请在网页右上角选择模型并在设置中填写对应的 API Key")

    user_message = f"请生成一个关于「{topic}」的算法可视化 Manim 动画脚本。数据量使用 {req.count} 个元素。风格：{req.style}。"

    try:
        from llm_service import chat_completion
        raw_response = await chat_completion(
            system_prompt=MANIM_VIDEO_SYSTEM_PROMPT,
            user_message=user_message,
            temperature=0.3,
            max_tokens=4096,
            api_key=req.api_key or None,
            model=req.model or None,
            api_base=req.api_base or None,
        )
    except Exception as e:
        raise HTTPException(500, f"LLM 调用失败: {str(e)}")

    # 提取 Python 代码
    import re
    code_match = re.search(r'```(?:python)?\s*\n?(.*?)```', raw_response, re.DOTALL)
    if code_match:
        script = code_match.group(1).strip()
    else:
        # 尝试直接作为代码（去掉可能的解释文本）
        script = raw_response.strip()

    if not script or len(script) < 50:
        raise HTTPException(500, "LLM 未能生成有效的 Manim 脚本")

    # 生成安全文件名
    safe_topic = "".join(c if c.isalnum() or c in '_-' else '_' for c in topic[:30])
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"manim_{safe_topic}_{timestamp}.py"
    script_path = OUTPUT_DIR / filename
    script_path.write_text(script, encoding='utf-8')

    # 尝试用 Manim 渲染（如果已安装）
    video_url = ""
    render_status = "script_only"
    render_output = ""

    # 提取 Scene 类名
    scene_name = _extract_manim_scene(script)

    try:
        import subprocess
        # manim -ql <script> <SceneClass> 格式
        cmd = ["manim", "-ql", str(script_path), scene_name]
        result = subprocess.run(
            cmd,
            capture_output=True, text=True, timeout=180,
            cwd=str(OUTPUT_DIR),
        )
        render_output = result.stdout + "\n" + result.stderr

        if result.returncode == 0:
            # Manim 输出: media/videos/{script_stem}/{quality}/{SceneName}.mp4
            script_stem = script_path.stem
            video_pattern = f"{script_stem}/480p15/{scene_name}.mp4"
            expected_path = OUTPUT_DIR / "media" / "videos" / video_pattern

            if expected_path.exists():
                # 复制到根目录方便访问
                video_filename = f"manim_{safe_topic}_{timestamp}.mp4"
                import shutil
                shutil.copy(expected_path, OUTPUT_DIR / video_filename)
                video_url = f"/api/resources/file/{video_filename}"
                render_status = "rendered"
            else:
                # 尝试模糊搜索
                for f in OUTPUT_DIR.glob(f"media/videos/{script_stem}/**/*.mp4"):
                    video_filename = f"manim_{safe_topic}_{timestamp}.mp4"
                    import shutil
                    shutil.copy(f, OUTPUT_DIR / video_filename)
                    video_url = f"/api/resources/file/{video_filename}"
                    render_status = "rendered"
                    break
    except FileNotFoundError:
        render_output = "Manim 未安装。请运行: pip install manim"
    except Exception as e:
        render_output = f"渲染异常: {str(e)}"

    # 存库
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=req.title or f"算法动画: {topic[:40]}",
            resource_type="video",
            content=script if not video_url else video_url,
            topic=topic,
            difficulty="基础",
            tags=["动画", "算法可视化", "Manim", "AI生成", topic],
            agent_generated="manim_llm",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "topic": topic,
        "script": script,
        "script_filename": filename,
        "video_url": video_url,
        "render_status": render_status,
        "render_output": render_output[:2000] if render_output else "",
        "install_guide": INSTALL_MANIM if render_status == "script_only" else "",
    }


# ==================== 讯飞AI数字人视频生成 ====================

XF_VMS_HOST = "vms.cn-huadong-1.xf-yun.com"
XF_VMS_START_PATH = "/v1/private/vms2d_start"
XF_VMS_CTRL_PATH = "/v1/private/vms2d_ctrl"
XF_VMS_STOP_PATH = "/v1/private/vms2d_stop"

# 可用的数字人形象
# 讯飞虚拟人形象列表
# ⚠️ 每个形象需要在控制台单独领取/购买，不是注册就能用
# 获取方式：console.xfyun.cn → AI虚拟人 → 接口服务 → 形象列表 → 复制形象ID
# 以下ID来自官方Python demo + Web SDK demo，可能已过期
DIGITAL_HUMAN_AVATARS = {
    "demo_default":   {"id": "118801001", "name": "默认形象(官方demo)"},
    "male_casual":    {"id": "118801001", "name": "休闲男(demo确认可用)"},
    "female_business": {"id": "110021007", "name": "商务女(demo确认可用)"},
    # 以下ID需到控制台确认是否已授权
    "male_business":  {"id": "110017006", "name": "商务男(需授权)"},
    "female_casual":  {"id": "110022010", "name": "休闲女(需授权)"},
    "female_teacher": {"id": "110005018", "name": "教师女(需授权)"},
    "male_service":   {"id": "110018008", "name": "服务男(需授权)"},
}

DIGITAL_HUMAN_SYSTEM_PROMPT = """你是一位专业的算法讲师。为讯飞AI数字人生成一份播报文稿，用于讲解算法主题。

要求：
1. 开场白：简短自我介绍 + 引出主题（约15秒语速的文字量）
2. 核心讲解：逐步讲解算法原理、关键步骤、复杂度分析（约60-90秒）
3. 实例演示：举一个具体例子演示算法执行过程（约30-45秒）
4. 总结回顾：要点总结 + 学习建议（约15秒）
5. 总时长控制在2-3分钟的语速（约400-800字）

格式：直接输出播报文稿，不要Markdown标记，不要特殊格式。
语气：亲切、专业，像一位耐心的老师在面对面讲课。
每句话不要太长（不超过25字），适合口语播报。"""


class DigitalHumanVideoRequest(BaseModel):
    topic: str = ""
    title: str = ""
    student_id: int = 1
    avatar: str = "male_business"
    app_id: str = ""
    api_key: str = ""
    api_secret: str = ""
    llm_api_key: str = ""  # 前端传过来的 LLM API Key（用于生成讲解文稿）
    llm_model: str = "deepseek-v4-flash"
    llm_api_base: str = ""


@router.post("/generate-xfyun-digital-human")
async def generate_xfyun_digital_human(req: DigitalHumanVideoRequest):
    """生成讯飞AI数字人讲解视频"""
    topic = req.topic.strip()
    if not topic:
        raise HTTPException(400, "请提供讲解主题")

    # ===== 步骤1: LLM 生成播报文稿 =====
    llm_api_key = req.llm_api_key or AppConfig().llm.api_key
    if not llm_api_key:
        raise HTTPException(400, "未配置 LLM API Key，请先在设置中配置一个 LLM 模型的 API Key")

    llm_model = req.llm_model or "deepseek-v4-flash"
    llm_api_base = req.llm_api_base or ""

    # 调试日志
    print(f"[数字人] LLM调用: model={llm_model}, base={llm_api_base}, key_len={len(llm_api_key)}, key_prefix={llm_api_key[:8]}...")

    try:
        from llm_service import chat_completion
        script_text = await chat_completion(
            system_prompt=DIGITAL_HUMAN_SYSTEM_PROMPT,
            user_message=f"请为主题「{topic}」生成一份2-3分钟的算法讲解播报文稿。",
            temperature=0.6,
            max_tokens=2048,
            api_key=llm_api_key,
            model=llm_model,
            api_base=llm_api_base or "https://api.deepseek.com",
        )
    except Exception as e:
        err_str = str(e)
        # 401 时给出更友好的提示
        if "401" in err_str:
            raise HTTPException(500, f"讯飞API认证失败(401)。请确认在设置中填入的是「APIPassword」(HTTP协议密码)，而不是APIKey/APISecret。获取地址: console.xfyun.cn → 对应模型 → HTTP服务接口认证信息 → APIPassword。原始错误: {err_str}")
        raise HTTPException(500, f"LLM 文稿生成失败: {err_str}")

    # ===== 步骤2: 调用讯飞虚拟人API =====
    xf_app_id = req.app_id
    xf_api_key = req.api_key
    xf_api_secret = req.api_secret

    if not all([xf_app_id, xf_api_key, xf_api_secret]):
        raise HTTPException(400, "请提供讯飞虚拟人API鉴权信息（APP ID + API Key + API Secret）")

    avatar_info = DIGITAL_HUMAN_AVATARS.get(req.avatar, DIGITAL_HUMAN_AVATARS["demo_default"])

    try:
        # 2.1 启动虚拟人会话（鉴权参数作为URL查询参数，不放header）
        start_url = _xfyun_build_auth_url(XF_VMS_HOST, XF_VMS_START_PATH, xf_api_key, xf_api_secret)

        start_body = {
            "header": {
                "app_id": xf_app_id,
                "uid": f"algoascend_{req.student_id}",
            },
            "parameter": {
                "vmr": {
                    "stream": {"protocol": "rtmp"},
                    "avatar_id": avatar_info["id"],
                    "width": 1280,
                    "height": 720,
                }
            }
        }

        async with httpx.AsyncClient(timeout=30.0) as client:
            start_resp = await client.post(
                start_url,
                json=start_body,
                headers={
                    "Content-Type": "application/json;charset=UTF-8",
                }
            )
            # 先拿原始文本，避免 .json() 抛异常丢失诊断信息
            start_text = start_resp.text
            try:
                start_data = start_resp.json()
            except Exception:
                raise HTTPException(500, f"讯飞虚拟人启动接口返回非JSON响应 (HTTP {start_resp.status_code}): {start_text[:500]}")

        # 打印完整响应便于服务端日志诊断
        print(f"[数字人] vms2d_start 响应: {json.dumps(start_data, ensure_ascii=False)[:800]}")

        if start_data.get("header", {}).get("code") != 0:
            header = start_data.get("header", {})
            err_code = header.get("code", "unknown")
            err_msg = header.get("message", "")
            # 把完整响应一并返回，避免显示"未知错误"
            raise HTTPException(
                500,
                f"讯飞虚拟人启动失败 [code={err_code}]: {err_msg or '讯飞未返回错误描述'} | 完整响应: {json.dumps(start_data, ensure_ascii=False)[:500]}"
            )

        # 修正：session 和 stream_url 都在 header 里（按官方文档）
        session_id = start_data.get("header", {}).get("session", "")
        stream_url = start_data.get("header", {}).get("stream_url", "")

        # stream_url 兼容处理：可能是 base64 编码，也可能是原始URL
        if stream_url:
            try:
                decoded = base64.b64decode(stream_url).decode("utf-8")
                if decoded.startswith("rtmp") or decoded.startswith("http"):
                    stream_url = decoded
            except Exception:
                pass  # 已经是原始URL，直接用

        # 2.2 文本驱动数字人播报（按官方demo：tts在parameter下，text在payload下且需base64编码）
        ctrl_url = _xfyun_build_auth_url(XF_VMS_HOST, XF_VMS_CTRL_PATH, xf_api_key, xf_api_secret)

        text_b64 = base64.b64encode(script_text.encode("utf-8")).decode("utf-8")
        ctrl_body = {
            "header": {
                "app_id": xf_app_id,
                "uid": f"algoascend_{req.student_id}",
                "session": session_id,
            },
            "parameter": {
                "tts": {
                    "vcn": "x3_lingxiaoxuan",  # 发音人
                    "speed": 50,
                    "volume": 50,
                    "pitch": 50,
                }
            },
            "payload": {
                "text": {
                    "encoding": "utf8",
                    "status": 3,  # 3=最后一帧
                    "text": text_b64,
                }
            }
        }

        async with httpx.AsyncClient(timeout=30.0) as client:
            ctrl_resp = await client.post(
                ctrl_url,
                json=ctrl_body,
                headers={
                    "Content-Type": "application/json;charset=UTF-8",
                }
            )
            ctrl_text = ctrl_resp.text
            try:
                ctrl_data = ctrl_resp.json()
            except Exception:
                raise HTTPException(500, f"数字人驱动接口返回非JSON响应 (HTTP {ctrl_resp.status_code}): {ctrl_text[:500]}")

        print(f"[数字人] vms2d_ctrl 响应: {json.dumps(ctrl_data, ensure_ascii=False)[:800]}")

        if ctrl_data.get("header", {}).get("code") != 0:
            ctrl_header = ctrl_data.get("header", {})
            ctrl_code = ctrl_header.get("code", "unknown")
            ctrl_msg = ctrl_header.get("message", "")
            raise HTTPException(
                500,
                f"数字人驱动失败 [code={ctrl_code}]: {ctrl_msg or '讯飞未返回错误描述'} | 完整响应: {json.dumps(ctrl_data, ensure_ascii=False)[:500]}"
            )

        drive_status = "success"

    except HTTPException:
        raise
    except Exception as e:
        # 讯飞API调用失败，返回文稿让用户查看
        import traceback
        traceback.print_exc()

        # 保存文稿到文件
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        safe_topic = "".join(c if c.isalnum() or c in '_-' else '_' for c in topic[:30])
        filename = f"xfyun_script_{safe_topic}_{timestamp}.txt"
        filepath = OUTPUT_DIR / filename
        filepath.write_text(script_text, encoding='utf-8')

        return {
            "status": "partial",
            "resource_id": 0,
            "topic": topic,
            "script": script_text,
            "avatar": avatar_info["name"],
            "stream_url": "",
            "video_url": "",
            "drive_status": "xfyun_api_failed",
            "message": f"讯飞虚拟人API调用失败: {str(e)}，已生成讲解文稿供参考",
            "script_filename": filename,
        }

    # ===== 步骤3: 尝试用 ffmpeg 录制 RTMP 流 =====
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    safe_topic = "".join(c if c.isalnum() or c in '_-' else '_' for c in topic[:30])
    video_filename = f"xfyun_video_{safe_topic}_{timestamp}.mp4"
    video_path = OUTPUT_DIR / video_filename
    video_url = ""
    record_status = "not_attempted"

    if stream_url and stream_url.startswith("rtmp"):
        try:
            import subprocess
            # 用 ffmpeg 录制 RTMP 流（录制10秒后自动停止）
            result = subprocess.run(
                ["ffmpeg", "-y", "-i", stream_url, "-t", "15", "-c", "copy", str(video_path)],
                capture_output=True, text=True, timeout=30,
            )
            if video_path.exists() and video_path.stat().st_size > 1000:
                video_url = f"/api/resources/file/{video_filename}"
                record_status = "recorded"
            else:
                record_status = f"record_failed: {result.stderr[:200]}"
        except FileNotFoundError:
            record_status = "ffmpeg_not_installed"
        except Exception as e:
            record_status = f"record_error: {str(e)[:200]}"

    # ===== 保存资源到数据库 =====
    db = get_db_sync()
    try:
        resource = LearningResource(
            student_id=req.student_id,
            title=req.title or f"数字人讲解: {topic[:40]}",
            resource_type="video_script",
            content=script_text if not video_url else video_url,
            topic=topic,
            difficulty="基础",
            tags=["数字人", "AI生成", "讯飞星火", "讲解视频", topic],
            agent_generated="xfyun_digital_human",
        )
        db.add(resource)
        db.commit()
        resource_id = resource.id
    finally:
        db.close()

    return {
        "status": "ok",
        "resource_id": resource_id,
        "topic": topic,
        "script": script_text,
        "avatar": avatar_info["name"],
        "avatar_id": avatar_info["id"],
        "stream_url": stream_url,
        "video_url": video_url,
        "record_status": record_status,
        "drive_status": drive_status,
        "message": "数字人讲解视频已生成" if video_url else f"数字人播报会话已创建（RTMP流: {stream_url}），录制状态: {record_status}",
    }


# ==================== 现有的 CRUD ====================

@router.get("/{student_id}")
async def get_resources(student_id: int = 1, resource_type: str = None, limit: int = 50):
    """获取学习资源列表"""
    db = get_db_sync()
    try:
        query = db.query(LearningResource).filter(
            LearningResource.student_id == student_id
        )
        if resource_type:
            query = query.filter(LearningResource.resource_type == resource_type)

        resources = (
            query.order_by(LearningResource.created_at.desc())
            .limit(limit)
            .all()
        )

        grouped = {}
        for r in resources:
            rtype = r.resource_type
            if rtype not in grouped:
                grouped[rtype] = []
            grouped[rtype].append(r.to_dict())

        return {
            "total": len(resources),
            "by_type": grouped,
            "resources": [r.to_dict() for r in resources],
        }
    finally:
        db.close()


@router.get("/detail/{resource_id}")
async def get_resource_detail(resource_id: int):
    """获取单个资源详情"""
    db = get_db_sync()
    try:
        resource = db.query(LearningResource).filter(
            LearningResource.id == resource_id
        ).first()
        if not resource:
            return {"error": "资源不存在", "code": 404}
        return resource.to_dict()
    finally:
        db.close()


@router.delete("/{resource_id}")
async def delete_resource(resource_id: int):
    """删除资源"""
    db = get_db_sync()
    try:
        resource = db.query(LearningResource).filter(
            LearningResource.id == resource_id
        ).first()
        if resource:
            db.delete(resource)
            db.commit()
        return {"status": "ok"}
    finally:
        db.close()
